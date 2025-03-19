from pptx import Presentation # module used to work with ppt's

def extract_txt_from_ppt(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text.strip()

file_path = "upload example.pptx"
extracted_text = extract_txt_from_ppt(file_path)
print(extracted_text)
